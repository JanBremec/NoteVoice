"""
faiss_rag.py

Production-ready skeleton for a FAISS-backed RAG document store.
Features:
 - extract text from PDF, PPTX, DOCX, TXT, MD
 - store documents in SQLite (with FTS5 for keyword search)
 - store vectors in FAISS (IndexFlatIP wrapped with IndexIDMap for persistent IDs)
 - pluggable embedding model (default: sentence-transformers)
 - functions: add_files, add_text, search_similarity, search_keyword, hybrid_search, save, load
 - chunking support for long documents
 - subject-based filtering

Dependencies (pip):
faiss-cpu, sentence-transformers, pdfplumber, python-docx, python-pptx, tqdm

Notes:
 - This is a single-file module intended as a starting point for production usage.
 - For heavy production workloads consider async ingestion, batching, multiprocessing, secure storage, encryption, monitoring.
"""

import os
import io
import json
import sqlite3
import math
from typing import List, Optional, Dict, Tuple, Iterable

import numpy as np
from tqdm import tqdm

# External libs
try:
    import faiss
except Exception as e:
    raise ImportError("faiss not installed. pip install faiss-cpu")

try:
    from sentence_transformers import SentenceTransformer
except Exception as e:
    SentenceTransformer = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


# --------------------------- Utilities ---------------------------

def ensure_dir(path: str):
    os.makedirs(os.path.dirname(path), exist_ok=True) if os.path.dirname(path) else None


def normalize_embeddings(embs: np.ndarray) -> np.ndarray:
    # normalize rows to unit length for cosine similarity (dot product)
    norms = np.linalg.norm(embs, axis=1, keepdims=True)
    norms[norms == 0] = 1.0
    return embs / norms

def chunk_text(text: str, size: int = 500, overlap: int = 50) -> List[str]:
    """
    Simple character-based chunking.
    """
    if not text:
        return []
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = start + size
        chunk = text[start:end]
        chunks.append(chunk)
        start += size - overlap
    return chunks

# --------------------------- Storage: SQLite with FTS5 ---------------------------

class DocumentStore:
    """
    SQLite-backed document store with FTS5 for keyword search.
    Stores: documents(id INTEGER PRIMARY KEY, content TEXT, metadata JSON string)
    Also stores an FTS virtual table 'docs_fts(content, metadata, id UNINDEXED)'
    """

    def __init__(self, db_path: str = "faiss_rag.db"):
        self.db_path = db_path
        self.conn = sqlite3.connect(self.db_path, check_same_thread=False)
        self._init_db()

    def _init_db(self):
        cur = self.conn.cursor()
        cur.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY,
                content TEXT NOT NULL,
                metadata TEXT
            );
        """)
        # FTS5 table for content + metadata search. Use direct insert into FTS for simplicity.
        cur.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS docs_fts USING fts5(content, metadata, tokenize = 'porter');
        """)
        self.conn.commit()

    def add(self, content: str, metadata: Optional[Dict] = None) -> int:
        cur = self.conn.cursor()
        metadata_json = json.dumps(metadata or {})
        cur.execute("INSERT INTO documents (content, metadata) VALUES (?, ?)", (content, metadata_json))
        doc_id = cur.lastrowid
        cur.execute("INSERT INTO docs_fts(rowid, content, metadata) VALUES (?, ?, ?)", (doc_id, content, metadata_json))
        self.conn.commit()
        return doc_id

    def get(self, doc_id: int) -> Optional[Dict]:
        cur = self.conn.cursor()
        cur.execute("SELECT id, content, metadata FROM documents WHERE id = ?", (doc_id,))
        row = cur.fetchone()
        if not row:
            return None
        return {"id": row[0], "content": row[1], "metadata": json.loads(row[2] or "{}")}

    def search_keyword(self, query: str, k: int = 10, subject: Optional[str] = None) -> List[Dict]:
        cur = self.conn.cursor()
        # FTS5 match
        cur.execute("SELECT rowid, content, metadata FROM docs_fts WHERE docs_fts MATCH ? LIMIT ?", (query, k * 5)) # Fetch more for filtering
        rows = cur.fetchall()
        results = []
        for r in rows:
            meta = json.loads(r[2] or "{}")
            if subject and meta.get('subject', '').lower() != subject.lower():
                continue
            results.append({"id": r[0], "content": r[1], "metadata": meta})
            if len(results) >= k:
                break
        return results

    def list_documents(self, subject: Optional[str] = None) -> List[Dict]:
        cur = self.conn.cursor()
        # We want to group by source_path or filename to avoid listing every chunk
        # But for now, let's just return all unique files based on metadata
        
        # Fetch all docs
        cur.execute("SELECT id, content, metadata FROM documents")
        rows = cur.fetchall()
        
        # Group by filename/source_path
        files_map = {}
        
        for r in rows:
            meta = json.loads(r[2] or "{}")
            
            # Filter by subject if provided
            if subject:
                doc_subject = meta.get('subject', '')
                if doc_subject.lower() != subject.lower():
                    continue
            
            filename = meta.get('filename', 'unknown')
            source_path = meta.get('source_path', filename)
            
            # Use source_path as key to identify unique files
            if source_path not in files_map:
                files_map[source_path] = {
                    "id": r[0], # Use ID of the first chunk found
                    "filename": filename,
                    "subject": meta.get('subject', 'Uncategorized'),
                    "type": os.path.splitext(filename)[1].replace('.', '') if filename != 'unknown' else 'txt',
                    "size": "Unknown", # We don't store size currently
                    "date": "Unknown", # We don't store date currently
                    "metadata": meta
                }
        
        return list(files_map.values())

    def get_subjects(self) -> List[str]:
        cur = self.conn.cursor()
        cur.execute("SELECT metadata FROM documents")
        rows = cur.fetchall()
        subjects = set()
        for r in rows:
            meta = json.loads(r[0] or "{}")
            if 'subject' in meta:
                subjects.add(meta['subject'])
        return list(subjects)

    def close(self):
        self.conn.commit()
        self.conn.close()


# --------------------------- FAISS Vector Store ---------------------------

class FaissVectorStore:
    """
    FAISS wrapper that persists index to disk and maps integer doc ids to FAISS ids via IndexIDMap.
    Uses inner IndexFlatIP (dot product) and expects normalized vectors for cosine.
    """

    def __init__(self, dim: int, index_path: str = "faiss.index"):
        self.dim = dim
        self.index_path = index_path
        self._init_index()

    def _init_index(self):
        if os.path.exists(self.index_path):
            try:
                self.load()
                return
            except Exception as e:
                print(f"Failed to load existing index: {e}. Creating new one.")
        
        quant = faiss.IndexFlatIP(self.dim)
        # IndexIDMap allows us to use our own ids (doc ids)
        self.index = faiss.IndexIDMap(quant)

    def add(self, vectors: np.ndarray, ids: np.ndarray):
        assert vectors.shape[1] == self.dim
        ids = ids.astype('int64')
        self.index.add_with_ids(vectors, ids)

    def search(self, vectors: np.ndarray, k: int = 10) -> Tuple[np.ndarray, np.ndarray]:
        # returns (scores, ids)
        scores, ids = self.index.search(vectors, k)
        return scores, ids

    def save(self, path: Optional[str] = None):
        path = path or self.index_path
        ensure_dir(path)
        faiss.write_index(self.index, path)

    def load(self, path: Optional[str] = None):
        path = path or self.index_path
        if not os.path.exists(path):
            raise FileNotFoundError(path)
        self.index = faiss.read_index(path)
        # make sure it's an IndexIDMap
        if not isinstance(self.index, faiss.IndexIDMap):
            raise RuntimeError("Loaded index is not IndexIDMap")


# --------------------------- Extractors ---------------------------

class Extractor:
    @staticmethod
    def from_file(path: str) -> str:
        ext = os.path.splitext(path)[1].lower()
        if ext in {'.txt', '.md'}:
            return Extractor._from_text(path)
        if ext in {'.pdf'}:
            return Extractor._from_pdf(path)
        if ext in {'.docx'}:
            return Extractor._from_docx(path)
        if ext in {'.pptx'}:
            return Extractor._from_pptx(path)
        raise ValueError(f"Unsupported file type: {ext}")

    @staticmethod
    def _from_text(path: str) -> str:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    @staticmethod
    def _from_pdf(path: str) -> str:
        if pdfplumber is None:
            raise ImportError('pdfplumber required to extract PDF text')
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                txt = page.extract_text()
                if txt:
                    text_parts.append(txt)
        return "\n".join(text_parts)

    @staticmethod
    def _from_docx(path: str) -> str:
        if docx is None:
            raise ImportError('python-docx required to extract DOCX text')
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paragraphs)

    @staticmethod
    def _from_pptx(path: str) -> str:
        if Presentation is None:
            raise ImportError('python-pptx required to extract PPTX text')
        prs = Presentation(path)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    t = shape.text
                    if t:
                        slides_text.append(t)
        return "\n".join(slides_text)


# --------------------------- Main RAG class ---------------------------

class FaissRAG:
    def __init__(self,
                 db_path: str = "faiss_rag.db",
                 index_path: str = "faiss.index",
                 embedding_model_name: str = "all-MiniLM-L6-v2",
                 batch_size: int = 64):
        if SentenceTransformer is None:
            raise ImportError('sentence-transformers is required for default embedding model')
        self.docstore = DocumentStore(db_path)
        self.embedder = SentenceTransformer(embedding_model_name)
        self.dim = self.embedder.get_sentence_embedding_dimension()
        self.vstore = FaissVectorStore(dim=self.dim, index_path=index_path)
        self.index_path = index_path
        self.batch_size = batch_size

    # ---------------- Ingestion ----------------
    def add_files(self, paths: Iterable[str], metadata: Optional[Dict] = None) -> List[int]:
        ids = []
        for p in paths:
            try:
                txt = Extractor.from_file(p)
            except Exception as e:
                print(f"Error extracting from {p}: {e}")
                continue
                
            chunks = chunk_text(txt)
            
            base_meta = metadata.copy() if metadata else {}
            base_meta['source_path'] = p
            if 'filename' not in base_meta:
                base_meta['filename'] = os.path.basename(p)
            
            for i, chunk in enumerate(chunks):
                meta = base_meta.copy()
                meta['chunk_index'] = i
                doc_id = self.docstore.add(chunk, meta)
                ids.append(doc_id)
                
        # embed and add newly added docs
        self._index_new_docs(ids)
        self.save() # Auto-save after adding
        return ids

    def add_text(self, text: str, metadata: Optional[Dict] = None) -> List[int]:
        chunks = chunk_text(text)
        ids = []
        for i, chunk in enumerate(chunks):
            meta = metadata.copy() if metadata else {}
            meta['chunk_index'] = i
            doc_id = self.docstore.add(chunk, meta)
            ids.append(doc_id)
            
        self._index_new_docs(ids)
        self.save() # Auto-save after adding
        return ids

    def _index_new_docs(self, doc_ids: List[int]):
        # fetch text for doc_ids, embed in batches, add to FAISS
        if not doc_ids:
            return
        batches = [doc_ids[i:i+self.batch_size] for i in range(0, len(doc_ids), self.batch_size)]
        for batch in batches:
            texts = [self.docstore.get(d)['content'] for d in batch]
            embeddings = self._embed_texts(texts)
            embeddings = normalize_embeddings(embeddings)
            ids = np.array(batch).astype('int64')
            self.vstore.add(embeddings, ids)

    def _embed_texts(self, texts: List[str]) -> np.ndarray:
        # sentence-transformers supports batching in encode
        embs = self.embedder.encode(texts, batch_size=self.batch_size, convert_to_numpy=True, show_progress_bar=False)
        if embs.ndim == 1:
            embs = np.expand_dims(embs, 0)
        return embs

    # ---------------- Persistence ----------------
    def save(self):
        self.vstore.save(self.index_path)

    def load(self):
        self.vstore.load(self.index_path)

    # ---------------- Search ----------------
    def search_similarity(self, query: str, k: int = 5, subject: Optional[str] = None) -> List[Dict]:
        q_emb = self._embed_texts([query])
        q_emb = normalize_embeddings(q_emb)
        
        # Fetch more candidates if filtering to avoid missing relevant docs
        # If subject is provided, we need to fetch enough to likely find matches
        fetch_k = max(k * 10, 100) if subject else k
        
        scores, ids = self.vstore.search(q_emb, fetch_k)
        
        results = []
        for i, doc_id in enumerate(ids[0]):
            if doc_id == -1:
                continue
            doc = self.docstore.get(int(doc_id))
            if not doc:
                continue
            
            # Filter by subject if provided
            if subject:
                doc_subject = doc['metadata'].get('subject', '')
                # Case-insensitive comparison
                if doc_subject.lower() != subject.lower():
                    continue
            
            results.append({
                "id": int(doc_id),
                "score": float(scores[0][i]),
                "content": doc['content'],
                "metadata": doc['metadata']
            })
            
            if len(results) >= k:
                break
                
        return results

    def search_keyword(self, query: str, k: int = 5, subject: Optional[str] = None) -> List[Dict]:
        # simple wrapper around sqlite FTS5
        return self.docstore.search_keyword(query, k, subject)

    def hybrid_search(self, query: str, k: int = 10, alpha: float = 0.5, subject: Optional[str] = None) -> List[Dict]:
        """
        Hybrid search: run both keyword (FTS) and vector similarity and merge results.
        alpha: weight for semantic score (0..1). Final score = alpha * semantic + (1-alpha) * keyword_score
        Keyword score is approximated using FTS rank ordering (we convert rank -> score)
        """
        semantic = self.search_similarity(query, k, subject=subject)
        keyword = self.search_keyword(query, k, subject=subject)
        
        # convert keyword list to dict of id->score (simple decreasing score)
        kw_scores = {}
        for rank, item in enumerate(keyword):
            # higher rank -> higher score
            kw_scores[item['id']] = 1.0 / (1 + rank)
            
        merged = {}
        for item in semantic:
            merged[int(item['id'])] = {'id': int(item['id']), 'semantic': item['score'], 'content': item['content'], 'metadata': item['metadata']}
            
        for item in keyword:
            did = int(item['id'])
            if did in merged:
                merged[did]['keyword'] = kw_scores[did]
            else:
                merged[did] = {'id': did, 'semantic': 0.0, 'keyword': kw_scores[did], 'content': item['content'], 'metadata': item['metadata']}
                
        results = []
        for v in merged.values():
            s_sem = v.get('semantic', 0.0)
            s_kw = v.get('keyword', 0.0)
            final = alpha * s_sem + (1 - alpha) * s_kw
            results.append({'id': v['id'], 'score': float(final), 'semantic': float(s_sem), 'keyword': float(s_kw), 'content': v['content'], 'metadata': v['metadata']})
            
        results = sorted(results, key=lambda x: x['score'], reverse=True)
        return results[:k]

    def list_documents(self, subject: Optional[str] = None) -> List[Dict]:
        return self.docstore.list_documents(subject)

    def get_subjects(self) -> List[str]:
        return self.docstore.get_subjects()

    def close(self):
        self.docstore.close()
